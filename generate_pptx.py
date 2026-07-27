"""
Green - Generador de presentaciones OpEx/DMAIC
------------------------------------------------
Llena el template Simplicity_Template_OpEx_DA_2026.pptx a partir de un JSON
con los datos del proyecto (ver schema/EJEMPLO.json).

Uso:
    python generate_pptx.py data/project.json output/Presentation.pptx
"""
import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
TEMPLATE_PATH = ROOT / "Simplicity_Template_OpEx_DA_2026.pptx"


# ---------------------------------------------------------------------------
# Helpers de bajo nivel
# ---------------------------------------------------------------------------

def shape_by_id(slide, shape_id):
    def search(shapes):
        for shape in shapes:
            if shape.shape_id == shape_id:
                return shape
            if shape.shape_type == 6:  # GROUP
                found = search(shape.shapes)
                if found is not None:
                    return found
        return None

    result = search(slide.shapes)
    if result is None:
        raise KeyError(f"shape_id {shape_id} no encontrado en la slide")
    return result


def set_run_value(paragraph, run_index, value):
    """Reemplaza solo el texto del run run_index dentro de un parrafo con
    varios runs (usado para pares 'Etiqueta: valor' donde la etiqueta debe
    conservarse en negritas)."""
    if run_index < len(paragraph.runs):
        paragraph.runs[run_index].text = str(value)


def set_single_text(shape, text):
    """Reemplaza el texto de un shape simple (1 parrafo, 1 o mas runs),
    conservando el formato del primer run."""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].text = text
    for extra in para.runs[1:]:
        extra.text = ""
    for p in tf.paragraphs[1:]:
        for r in p.runs:
            r.text = ""


def set_cell_text(cell, text):
    tf = cell.text_frame
    para = tf.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].text = str(text)
    for extra in para.runs[1:]:
        extra.text = ""
    for p in tf.paragraphs[1:]:
        for r in p.runs:
            r.text = ""


def find_header_index(paragraphs, header_substr):
    for i, p in enumerate(paragraphs):
        if header_substr.lower() in p.text.lower():
            return i
    return None


def replace_bullet_group(text_frame, header_substr, items):
    """Sustituye el grupo de bullets que sigue a un parrafo 'header_substr'
    hasta el siguiente parrafo vacio (o fin del cuadro), por 'items'."""
    paragraphs = text_frame.paragraphs
    start = find_header_index(paragraphs, header_substr)
    if start is None:
        return
    start += 1
    end = start
    while end < len(paragraphs) and paragraphs[end].text.strip() != "":
        end += 1

    if start >= len(paragraphs):
        return

    template_p = paragraphs[start]._p
    parent = template_p.getparent()
    ref_element = paragraphs[end]._p if end < len(paragraphs) else None

    tmpl_clone_source = copy.deepcopy(template_p)

    # quitar los parrafos originales del grupo
    for idx in range(start, end):
        el = paragraphs[idx]._p
        parent.remove(el)

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for item in items:
        new_p = copy.deepcopy(tmpl_clone_source)
        runs = new_p.findall("a:r", ns)
        if runs:
            first = runs[0]
            t = first.find("a:t", ns)
            if t is None:
                t = first.makeelement(f"{{{ns['a']}}}t", {})
                first.append(t)
            t.text = str(item)
            for extra in runs[1:]:
                new_p.remove(extra)
        if ref_element is not None:
            ref_element.addprevious(new_p)
        else:
            parent.append(new_p)


def set_paragraph_text(text_frame, para_index, text):
    tf_paras = text_frame.paragraphs
    if para_index >= len(tf_paras):
        return
    para = tf_paras[para_index]
    if not para.runs:
        para.add_run()
    para.runs[0].text = text
    for r in para.runs[1:]:
        r.text = ""


def fill_table_row(table, row_idx, values, start_col=0):
    row = table.rows[row_idx]
    for i, value in enumerate(values):
        set_cell_text(row.cells[start_col + i], value)


# ---------------------------------------------------------------------------
# Llenado por slide
# ---------------------------------------------------------------------------

def fill_slide1(slide, d):
    define = d.get("define", {})
    set_single_text(shape_by_id(slide, 21), d.get("project_title", ""))

    tb107 = shape_by_id(slide, 107)
    set_paragraph_text(tb107.text_frame, 1, define.get("problem_statement", ""))
    replace_bullet_group(tb107.text_frame, "Project Objective", define.get("objective_bullets", []))
    replace_bullet_group(tb107.text_frame, "Project Benefits", define.get("benefit_bullets", []))

    if define.get("glossary_text"):
        set_single_text(shape_by_id(slide, 23), define["glossary_text"])

    # Roles y fechas (viven dentro de "Group 11", id=12)
    leader_sponsor = shape_by_id(slide, 108)
    set_run_value(leader_sponsor.text_frame.paragraphs[0], 2, define.get("project_leader", ""))
    set_run_value(leader_sponsor.text_frame.paragraphs[0], 4, define.get("sponsor", ""))
    set_run_value(leader_sponsor.text_frame.paragraphs[1], 1, define.get("team_members", ""))

    opex_financial = shape_by_id(slide, 3)
    set_run_value(opex_financial.text_frame.paragraphs[0], 1, define.get("opex_master", ""))
    set_run_value(opex_financial.text_frame.paragraphs[1], 1, define.get("financial_approver", ""))

    dates = shape_by_id(slide, 66)
    set_run_value(dates.text_frame.paragraphs[0], 1, define.get("start_date", ""))
    set_run_value(dates.text_frame.paragraphs[0], 3, define.get("end_date", ""))

    # Executive champion / business unit approver (nombres sueltos, fuera del grupo)
    set_single_text(shape_by_id(slide, 27), define.get("executive_champion", ""))
    set_single_text(shape_by_id(slide, 29), define.get("business_unit_approver", ""))


def fill_slide2(slide, d):
    ib = d.get("intangible_benefits", {})
    table = shape_by_id(slide, 33).table

    def pad(lst, length, fill=""):
        lst = list(lst) if lst else []
        return (lst + [fill] * length)[:length]

    steps = pad(ib.get("process_steps"), 3, "")
    fill_table_row(table, 2, [steps[0], "", steps[1], "", steps[2], "", "Months per year", "", "Annual hours worked"])
    fill_table_row(table, 8, [steps[0], "", steps[1], "", steps[2], "", "Months per year", "", "Annual hours worked"])

    av = pad(ib.get("actual_values"), 5, 0)
    fill_table_row(table, 4, [av[0], "+", av[1], "+", av[2], "x", av[3], "=", av[4]])
    ov = pad(ib.get("optimized_values"), 5, 0)
    fill_table_row(table, 9, [ov[0], "+", ov[1], "+", ov[2], "x", ov[3], "=", ov[4]])

    if ib.get("before_list"):
        replace_free_list(shape_by_id(slide, 35), ib["before_list"])
    if ib.get("after_list"):
        replace_free_list(shape_by_id(slide, 40), ib["after_list"])

    if ib.get("annual_saving_text"):
        set_single_text(shape_by_id(slide, 34), ib["annual_saving_text"])
    if ib.get("confirmation_note"):
        set_single_text(shape_by_id(slide, 41), ib["confirmation_note"])


def replace_free_list(shape, items):
    """Para shapes tipo 'Title 3' donde cada parrafo ya es un bullet (sin header)."""
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    if not paragraphs:
        return
    template_p = paragraphs[0]._p
    parent = template_p.getparent()
    tmpl_clone_source = copy.deepcopy(template_p)
    for p in list(paragraphs):
        parent.remove(p._p)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for item in items:
        new_p = copy.deepcopy(tmpl_clone_source)
        runs = new_p.findall("a:r", ns)
        if runs:
            t = runs[0].find("a:t", ns)
            if t is not None:
                t.text = str(item)
            for extra in runs[1:]:
                new_p.remove(extra)
        parent.append(new_p)


def to_number(value, default=0):
    """Convierte a numero de forma segura (el agente a veces manda '48' en
    vez de 48, o texto vacio); si no se puede, regresa el default."""
    if isinstance(value, (int, float)):
        return value
    try:
        cleaned = str(value).replace(",", "").replace("$", "").replace("%", "").strip()
        return float(cleaned) if cleaned else default
    except (ValueError, TypeError):
        return default


def fill_slide3(slide, d):
    tb = d.get("tangible_benefits", {})
    cost = to_number(tb.get("cost_per_lost_project", 0))
    lost_current = to_number(tb.get("lost_projects_current", 0))
    lost_opt = to_number(tb.get("lost_projects_optimized", 0))
    total_current = cost * lost_current
    total_opt = cost * lost_opt

    set_single_text(shape_by_id(slide, 7),
                     f"${cost:,.0f} USD   x    {lost_current}    x   =   ${total_current:,.0f} USD")
    set_single_text(shape_by_id(slide, 4),
                     f"${cost:,.0f} USD   x    {lost_opt}    x   =   ${total_opt:,.2f} USD")

    if tb.get("summary_lines"):
        summary_shape = shape_by_id(slide, 5)
        tf = summary_shape.text_frame
        for i, line in enumerate(tb["summary_lines"][:len(tf.paragraphs)]):
            set_paragraph_text(tf, i, line)
    if tb.get("projected_note"):
        note_shape = shape_by_id(slide, 6)
        set_single_text(note_shape, tb["projected_note"])


def fill_slide4(slide, d):
    rows = d.get("implementation_plan", [])
    table = shape_by_id(slide, 24).table
    max_rows = 13
    for i in range(max_rows):
        row_idx = 2 + i
        if i < len(rows):
            r = rows[i]
            fill_table_row(table, row_idx, [
                str(i + 1),
                r.get("action", ""),
                r.get("responsible", ""),
                r.get("start_date", ""),
                r.get("completion_date", ""),
                r.get("status", ""),
                r.get("phase", ""),
                r.get("comments", ""),
            ])
        else:
            fill_table_row(table, row_idx, [str(i + 1), "", "", "", "", "", "", ""])


def replace_picture(shape, image_path):
    """Sustituye la imagen de un shape tipo Picture, conservando su posicion
    y tamano originales en la slide."""
    slide_part = shape.part
    image_part, rId = slide_part.get_or_add_image_part(str(image_path))
    for blip in shape._element.findall(".//" + qn("a:blip")):
        blip.set(qn("r:embed"), rId)


# slide 5: mapa de "slot logico" -> shape_id de la imagen de evidencia
IMAGE_SLOTS_SLIDE5 = {
    "before1": 3,
    "after1": 10,
    "before2": 4,
    "after2": 5,
}


def fill_slide5_images(slide, d):
    images = d.get("images", {})
    for slot, shape_id in IMAGE_SLOTS_SLIDE5.items():
        rel_path = images.get(slot)
        if not rel_path:
            continue
        img_path = ROOT / rel_path
        if not img_path.exists():
            print(f"AVISO: no se encontro la imagen {img_path}, se deja la original")
            continue
        try:
            replace_picture(shape_by_id(slide, shape_id), img_path)
        except Exception as e:
            print(f"AVISO: no se pudo reemplazar imagen '{slot}': {e}")


def fill_slide5(slide, d):
    se = d.get("solutions_evidence", {})
    if se.get("before_text"):
        set_single_text(shape_by_id(slide, 2), se["before_text"])
    if se.get("after_text"):
        set_single_text(shape_by_id(slide, 20), se["after_text"])
    if se.get("pilot_note_1"):
        set_single_text(shape_by_id(slide, 8), se["pilot_note_1"])
    if se.get("pilot_note_2"):
        set_single_text(shape_by_id(slide, 11), se["pilot_note_2"])
    if se.get("process_name"):
        set_single_text(shape_by_id(slide, 101), se["process_name"])
    if se.get("generate_reports_label"):
        set_single_text(shape_by_id(slide, 16), se["generate_reports_label"])


def fill_slide6(slide, d):
    rows = d.get("control_plan", [])
    table = shape_by_id(slide, 26).table
    max_rows = 4
    for i in range(max_rows):
        row_idx = 1 + i
        if i < len(rows):
            r = rows[i]
            fill_table_row(table, row_idx, [
                r.get("x", ""),
                r.get("target", ""),
                r.get("method", ""),
                r.get("action", ""),
                r.get("frequency", "Monthly"),
                r.get("responsible", ""),
                r.get("start_date", ""),
                r.get("reaction_plan", ""),
            ])
        else:
            fill_table_row(table, row_idx, ["", "", "", "", "", "", "", ""])


def fill_slide7(slide, d):
    mp = d.get("monitoring_plan")
    if not mp:
        return
    metric1_ids = [45, 56, 58, 60, 62]  # 5 valores
    metric2_ids = [54, 64, 66, 68, 70]  # 5 valores
    v1 = mp.get("metric1_values", [])
    v2 = mp.get("metric2_values", [])
    for shape_id, value in zip(metric1_ids, v1):
        set_single_text(shape_by_id(slide, shape_id), str(value))
    for shape_id, value in zip(metric2_ids, v2):
        set_single_text(shape_by_id(slide, shape_id), str(value))
    if mp.get("year_label"):
        for sid in (10, 25, 26):
            set_single_text(shape_by_id(slide, sid), mp["year_label"])


CHECKLIST_ORDER = [
    "Standard Tool",
    "Access Control",
    "Automated Feed*",
    "User & Back Up // Admin & Tool Maintenance Backup",
    "Documented SOP",
    "Hosted on Server*",
    "Training",
]


def fill_slide8(slide, d):
    tg = d.get("tollgate", {})
    table = shape_by_id(slide, 10).table
    checklist = tg.get("checklist", [])
    for i in range(7):
        row_idx = 1 + i
        criteria_label = CHECKLIST_ORDER[i]
        if i < len(checklist):
            item = checklist[i]
            set_cell_text(table.rows[row_idx].cells[0], item.get("criteria", criteria_label))
            set_cell_text(table.rows[row_idx].cells[1], item.get("yn", ""))
            set_cell_text(table.rows[row_idx].cells[2], item.get("comments", ""))

    if tg.get("date_approved"):
        header_shape = shape_by_id(slide, 13)
        tf = header_shape.text_frame
        set_paragraph_text(tf, 0, f"CONTROL\t\t Date Tollgate Approved: {tg['date_approved']}")

    if tg.get("comments_conclusion"):
        cc_shape = shape_by_id(slide, 12)
        tf = cc_shape.text_frame
        lines = tg["comments_conclusion"]
        for i, line in enumerate(lines[: max(0, len(tf.paragraphs) - 1)]):
            set_paragraph_text(tf, i + 1, line)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build(data, output_path):
    prs = Presentation(str(TEMPLATE_PATH))
    slides = list(prs.slides)

    steps = [
        ("slide 1 (Define)", fill_slide1, slides[0]),
        ("slide 2 (Intangible Benefits)", fill_slide2, slides[1]),
        ("slide 3 (Tangible Benefits)", fill_slide3, slides[2]),
        ("slide 4 (Implementation Plan)", fill_slide4, slides[3]),
        ("slide 5 (Solutions Evidence)", fill_slide5, slides[4]),
        ("slide 5 (imagenes)", fill_slide5_images, slides[4]),
        ("slide 6 (Control Plan)", fill_slide6, slides[5]),
        ("slide 7 (Monitoring Plan)", fill_slide7, slides[6]),
        ("slide 8 (Tollgate)", fill_slide8, slides[7]),
    ]
    errors = []
    for label, fn, slide in steps:
        try:
            fn(slide, data)
        except Exception as e:
            errors.append(f"{label}: {type(e).__name__}: {e}")
            print(f"AVISO: fallo en {label}, se deja el contenido original de esa parte. Detalle: {e}")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"OK -> {output_path}")
    if errors:
        print("\nATENCION: estas partes no se pudieron llenar (revisa el JSON):")
        for e in errors:
            print(f"  - {e}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Uso: python generate_pptx.py <data.json> <output.pptx>")
        sys.exit(1)
    data_path, out_path = sys.argv[1], sys.argv[2]
    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    build(data, out_path)
